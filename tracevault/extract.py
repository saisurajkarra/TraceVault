"""Make every file 'speak': translate any file's bytes into language.

The goal is that EVERY file — not just code/text — has a meaningful language
representation in the vector DB, so it is semantically addressable and can be
connected to what people and the AI said about it. Translation by type:

  text / code        -> the text itself (snippet)
  images             -> a natural-language caption from a local vision model (BLIP)
  PDF / DOCX / PPTX  -> extracted document text
  notebooks (.ipynb) -> markdown + code cell sources
  other binaries     -> typed metadata (kind, size)

The image model is loaded locally and lazily (only when an image is first seen).
If it cannot be loaded, this raises loudly — multimodal ingestion has no silent
fallback. Per-file extraction errors are logged and degrade to a typed stub for
that one file (the file still appears; nothing is fabricated).
"""

from __future__ import annotations

import io
import json
import logging
import os
from dataclasses import dataclass
from typing import Any

logger = logging.getLogger(__name__)

SNIPPET_CHARS = 4000
DOC_CHARS = 6000
CAPTION_MAX_BYTES = 12 * 1024 * 1024  # skip captioning images larger than this

IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff", ".tif"}
CODE_EXTS = {
    ".py", ".ts", ".tsx", ".js", ".jsx", ".go", ".rs", ".java", ".kt", ".c", ".h", ".cpp",
    ".cc", ".cs", ".rb", ".php", ".swift", ".scala", ".sh", ".sql", ".css", ".scss",
}

IMAGE_CAPTION_MODEL = "Salesforce/blip-image-captioning-base"


class ExtractError(RuntimeError):
    """Raised when a required extraction model cannot be loaded."""


@dataclass
class Extracted:
    text: str  # the language representation — what the file "says"
    modality: str  # text | code | image | pdf | docx | pptx | notebook | binary
    is_binary: bool


# --- lazy local image captioner (BLIP, loaded directly: pipeline task names drift) ---

_blip_processor: Any = None
_blip_model: Any = None


def _get_blip() -> tuple[Any, Any]:
    global _blip_processor, _blip_model
    if _blip_model is None:
        try:
            from transformers import (  # noqa: PLC0415  (lazy, heavy import)
                BlipForConditionalGeneration,
                BlipProcessor,
            )

            logger.info("Loading local image captioning model %r ...", IMAGE_CAPTION_MODEL)
            _blip_processor = BlipProcessor.from_pretrained(IMAGE_CAPTION_MODEL, use_fast=True)
            _blip_model = BlipForConditionalGeneration.from_pretrained(IMAGE_CAPTION_MODEL)
        except Exception as exc:  # model/download/dep failure
            raise ExtractError(
                "Could not load the local image-captioning model (BLIP). Multimodal ingest "
                "requires it and there is no fallback. The first run downloads the model; "
                f"check your network or model cache. Underlying error: {exc}"
            ) from exc
    return _blip_processor, _blip_model


def _caption_image(data: bytes, path: str) -> str:
    try:
        import torch  # noqa: PLC0415
        from PIL import Image  # noqa: PLC0415

        img = Image.open(io.BytesIO(data)).convert("RGB")
        processor, model = _get_blip()
        inputs = processor(images=img, return_tensors="pt")
        with torch.no_grad():
            out = model.generate(**inputs, max_new_tokens=40)
        caption: str = processor.decode(out[0], skip_special_tokens=True).strip()
        return caption or "image"
    except ExtractError:
        raise
    except Exception as exc:  # one unreadable image shouldn't abort the run
        logger.warning("Could not caption image %s: %s", path, exc)
        return "image (uncaptioned)"


def caption_images(datas: list[bytes], paths: list[str], *, batch_size: int = 8) -> list[str]:
    """Batch-caption images locally (much faster than one-by-one). Returns one caption per input.

    Loads BLIP once (raising loudly if it cannot load) and runs batched generation. Individual
    undecodable images degrade to a stub and are logged; the batch still completes.
    """
    if not datas:
        return []
    import torch  # noqa: PLC0415
    from PIL import Image  # noqa: PLC0415

    processor, model = _get_blip()
    out_captions: list[str] = []
    for start in range(0, len(datas), batch_size):
        chunk = datas[start : start + batch_size]
        chunk_paths = paths[start : start + batch_size]
        imgs: list[Any] = []
        idx_map: list[int] = []
        results = ["image (uncaptioned)"] * len(chunk)
        for i, d in enumerate(chunk):
            try:
                imgs.append(Image.open(io.BytesIO(d)).convert("RGB"))
                idx_map.append(i)
            except Exception as exc:
                logger.warning("Could not open image %s: %s", chunk_paths[i], exc)
        if imgs:
            inputs = processor(images=imgs, return_tensors="pt")
            with torch.no_grad():
                gen = model.generate(**inputs, max_new_tokens=40)
            decoded = [processor.decode(g, skip_special_tokens=True).strip() or "image" for g in gen]
            for j, orig_i in enumerate(idx_map):
                results[orig_i] = decoded[j]
        out_captions.extend(results)
    return out_captions


def _pdf_text(data: bytes) -> str:
    from pypdf import PdfReader  # noqa: PLC0415

    reader = PdfReader(io.BytesIO(data))
    parts = []
    for page in reader.pages[:30]:
        parts.append(page.extract_text() or "")
        if sum(len(p) for p in parts) > DOC_CHARS:
            break
    return "\n".join(parts)[:DOC_CHARS]


def _docx_text(data: bytes) -> str:
    import docx  # noqa: PLC0415

    document = docx.Document(io.BytesIO(data))
    return "\n".join(p.text for p in document.paragraphs if p.text)[:DOC_CHARS]


def _pptx_text(data: bytes) -> str:
    from pptx import Presentation  # noqa: PLC0415

    prs = Presentation(io.BytesIO(data))
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                lines.append(shape.text_frame.text)
    return "\n".join(t for t in lines if t)[:DOC_CHARS]


def _notebook_text(data: bytes) -> str:
    nb = json.loads(data.decode("utf-8", errors="replace"))
    parts = []
    for cell in nb.get("cells", []):
        src = cell.get("source", "")
        parts.append("".join(src) if isinstance(src, list) else str(src))
    return "\n".join(parts)[:DOC_CHARS]


def extract(path: str, data: bytes, *, enable_images: bool = True) -> Extracted:
    """Translate a file's bytes into a language representation it 'speaks'."""
    ext = os.path.splitext(path)[1].lower()

    if ext in IMAGE_EXTS:
        if enable_images and 0 < len(data) <= CAPTION_MAX_BYTES:
            return Extracted(text=f"[image] {_caption_image(data, path)}", modality="image", is_binary=True)
        return Extracted(text=f"[image file {ext.lstrip('.')}]", modality="image", is_binary=True)

    try:
        if ext == ".pdf":
            return Extracted(text=_pdf_text(data), modality="pdf", is_binary=True)
        if ext == ".docx":
            return Extracted(text=_docx_text(data), modality="docx", is_binary=True)
        if ext == ".pptx":
            return Extracted(text=_pptx_text(data), modality="pptx", is_binary=True)
        if ext == ".ipynb":
            return Extracted(text=_notebook_text(data), modality="notebook", is_binary=False)
    except Exception as exc:  # corrupt/unsupported doc — degrade to a typed stub, logged
        logger.warning("Could not extract %s (%s): %s", path, ext, exc)
        return Extracted(text=f"[{ext.lstrip('.')} document, {len(data)} bytes]", modality="binary", is_binary=True)

    # text / code: decode as UTF-8 unless it looks binary.
    if b"\x00" not in data[:8000]:
        try:
            text = data.decode("utf-8")
            modality = "code" if ext in CODE_EXTS else "text"
            return Extracted(text=text[:SNIPPET_CHARS], modality=modality, is_binary=False)
        except UnicodeDecodeError:
            pass

    return Extracted(
        text=f"[binary file {ext.lstrip('.') or 'no-ext'}, {len(data)} bytes]",
        modality="binary",
        is_binary=True,
    )
